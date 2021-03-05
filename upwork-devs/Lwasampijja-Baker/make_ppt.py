import pandas as pd
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.oxml.xmlchemy import OxmlElement

d = ''

# Glasswall palette
dark_blue                                   = RGBColor(14, 61, 90)
green_blue                                  = RGBColor(26, 145, 154)
white                                       = RGBColor(255, 255, 255)
blue1                                       = RGBColor(22, 94, 122)
# table colors:
gray                                        = RGBColor(191, 191, 191)
blue2                                       = RGBColor(45, 92, 117)
# rag colors:
green                                       = RGBColor(0, 204, 153)  # (0, 255, 0)
amber                                       = RGBColor(255, 204, 0)  # (255, 153, 51)
red                                         = RGBColor(255, 102, 102)  # (255, 0, 0)

# Letter font
gw_font                                     = 'Arial'


def examine_template():
    """
    Print default Master's slides information, title, subtitle, placeholders, etc.
    """
    prs                                    = Presentation()

    for n in range(0, 11):

        slide                              = prs.slides.add_slide(prs.slide_layouts[n])

        print('Master Slide ' + str(n))

        for shape in slide.placeholders:
            print('%d, %s' % (shape.placeholder_format.idx, shape.name))


def logo(slide, img_path=d + 'gw.png', place='top right'):
    """
    Insert logo in slide.
    :param slide: slide from presentation
    :param img_path: path to image file, Glasswall logo for default
    :param place: place to locate image, top right, center or top left
    """
    if place                               == 'top right':
        # Logo size
        width                              = Inches(1.2)
        height                             = Inches(0.6)  # width half
        top                                = Inches(0.1)
        left                               = Inches(10.0) - width - Inches(0.2)

    elif place                             == 'center':
        width                              = Inches(6.0)
        height                             = Inches(3.0)  # width halfs
        left                               = (Inches(10.0) - width) / 2
        top                                = (Inches(7.5) - height) / 2

    elif place                             == 'top left':
        width                              = Inches(1.25)
        height                             = Inches(0.625)  # width half
        top                                = Inches(0.25)
        left                               = Inches(0.3)

    pic                                    = slide.shapes.add_picture(img_path, left, top, width, height)


def set_background_color(slide, bg_color=dark_blue):
    """
    Set slide background color.
    :param slide: slide from presentation
    :param bg_color: background color
    """
    background                             = slide.background
    fill(background, bg_color)


def fill(shape, fill_color=dark_blue):
    """
    Fill shape with color.

    :param shape: MSO_SHAPE shape (MSO_SHAPE.RECTANGLE, MSO_SHAPE.ELLIPSE, etc)
    :param fill_color: fill color
    """
    fill                                   = shape.fill
    fill.solid()
    fill.fore_color.rgb                    = fill_color

def wrap_by_word(s, n):
    """
    Returns a string where \n is inserted between every n words

    :param s: string
    :param n: integer, number of words
    :return:
    """
    a                                      = s.split()
    ret                                    = ''
    for i in range(0, len(a), n):
        ret += ' '.join(a[i:i+n]) + '\n'
    return ret


def wrap_by_char(s, n):
    return '\n'.join(l for line in s.splitlines() for l in textwrap.wrap(line, width=n))

def get_data():
    # Create URL to JSON file (alternatively this can be a filepath)
    url                                    = 'https://wmwaredata.s3.us-east-2.amazonaws.com/gw_releases.json'
    # Load the first sheet of the JSON file into a data frame
    df                                     = pd.read_json(url, orient='columns')
    df                                     = df.rename(columns={'sub_repo_commit_url': 'sub_repo_url'})
    repos                                  = []
    dates                                  = []
    tags                                   = []
    hashes                                 = []
    descriptions                           = []

    for i in range(len(df)):
        # Repo
        repo                               = df['repo_name'].iloc[i] + '\n\n' + df['repo_url'].iloc[i]
        repos.append(repo)
        # Date
        d                                  = df['release_date'].iloc[i]
        if d is not None:
            d                              = d.split('T')
            date                           = d[0] + '\n\n' + d[1][:-1]
        else:
            date                           = ''
        dates.append(date)
        # Version / Tag
        t                                  = df['version'].iloc[i]
        tags.append(t)
        # Hash
        h                                  = df['hash'].iloc[i]
        hashes.append(h)
        # Notes / Description
        content                            = re.sub('<.*?>', '', df['release_notes'].iloc[i])
        des                                = wrap_by_word(content, n=20)
        descriptions.append(des)

        # Sub Repo
        s_name                             = df['sub_repo_name'].iloc[i]
        if s_name is not None:
            s_repo                         = s_name + '\n\n' + df['sub_repo_url'].iloc[i]
            repos.append(s_repo)
            # date
            dates.append(date)
            # tag
            tags.append(t)
            # Sub Hash
            s_h                            = df['sub_hash'].iloc[i]
            if s_h is not None:
                hashes.append(s_h)
            # notes
            descriptions.append(des)

    df                                     = pd.DataFrame()
    df['Repo']                             = repos
    df['Date']                             = dates
    df['Version']                          = tags
    df['Hash']                             = hashes
    df['Notes']                            = descriptions


    # Sort columns
    df                                     = df[['Repo', 'Date', 'Version', 'Hash', 'Notes']]

    # drop repeated rows
    df1                                    = df.drop_duplicates().reset_index(drop=True)
    return df1

def make_presentation(output_to, single=False, dm=False):
    """
    Autocreate power point presentation for 'Project Team Structure' sheet.

    :param sheet_name: sheet name, in this case 'Projects Team Structure'
    :param output_to: path to save output pptx file
    :param single: bool, create a single presentations and save it to single folder if true
    :param dm: bool, create a presentation per delivery manager and save it to dm folder if true
    """
    prs                                     = Presentation()
    df1                                     = get_data()

    # PROJECT SLIDES
    for row_index in range(len(df1)):
        add_project_slide(prs, df1, row_index)
    prs.save(output_to)

def add_project_slide(prs, df, row_index):
    """
    Add slide to presentation.

    :param prs: presentation
    :param df: pandas dataframe with presentation information
    :param row_index: index of the row with the information corresponding to the slide
    """
    repo                                    = df.iloc[row_index]['Repo']
    date                                    = df.iloc[row_index]['Date']
    version                                 = df.iloc[row_index]['Version']
    hash                                    = df.iloc[row_index]['Hash']
    notes                                   = df.iloc[row_index]['Notes']

    title_only_slide_layout                 = prs.slide_layouts[5]
    slide                                   = prs.slides.add_slide(title_only_slide_layout)
    title                                   = slide.shapes.title.text = "GW Releases"

    #set_background_color(slide)

    logo(slide)

    shapes                                  = slide.shapes

    # TITLE
    #title = shapes.title


    #title = '\n GW Releases' + repo.upper() + '\n'
    #text_settings(title, i=0)
    #text_settings(title, i=1)
    #text_settings(title, i=2, font_size=Pt(26))


    #text_settings(title, i=4, font_size=Pt(24), font_color=green_blue)

    rnr                                     = df[(df['Hash'] == hash)].reset_index()
    if len(rnr) > 0:
        add_table(shapes, rnr, blue1)


def text_settings(
        shape,
        i=0,
        alignment=PP_ALIGN.LEFT,
        font_color=white,
        font_size=Pt(9),
        font= gw_font,
        bold=False):
    """
    Format shape's text with alignment, font, font color and size, etc.

    :param shape: MSO_SHAPE shape (MSO_SHAPE.RECTANGLE, MSO_SHAPE.ELLIPSE, etc)
    :param i: line position
    :param alignment: alignment (PP_ALIGN.LEFT, PP_ALIGN.CENTER, etc.)
    :param font_color: font color
    :param font_size: font size
    :param font: letter font
    :param bold: bool, use bold letters if true
    """
    text                                       = shape.text_frame.paragraphs[i]
    text.alignment                             = alignment
    text.font.name                             = font
    text.font.size                             = font_size
    text.font.color.rgb                        = font_color
    text.font.bold                             = bold
    for paragraph in shape.text_frame.paragraphs:
        paragraph.font.size                    = Pt(9)
        paragraph.font.color.rgb               = RGBColor(255, 255, 255)


def add_table(
        shapes,
        df,
        table_color,
        top=Inches(1.5),
        col_width=Inches(3.0),
        left=Inches(0.3),
        width=Inches(3.5),
        height=Inches(0.5)):
    """
    Add table to slide.

    :param shapes: shapes attribute from slide (which in turn is an attribute of the presentation)
    :param df: pandas dataframe with 'Resource' and 'Responsability' information in columns
    :param table_color: table color
    :param top: distance (in inches) to top edge of slide (each slide is 10 per 7.5 inches)
    :param col_width: column width
    :param left: distance (in inches) to left edge of slide
    :param width: table width
    :param height: table height
    """
    cols = 5
    rows = len(df) + 1

    shape = shapes.add_table(rows, cols, left, top, width, height)
    table = shape.table

    # set column widths
    table.columns[0].width                     = Inches(1.5)
    table.columns[1].width                     = Inches(1.0)
    table.columns[2].width                     = Inches(1.0)
    table.columns[3].width                     = col_width
    table.columns[4].width                     = col_width

    # write column headings
    table.cell(0, 0).text                      = 'Repo'.capitalize()
    table.cell(0, 1).text                      = 'Date'.capitalize()
    table.cell(0, 2).text                      = 'Version'.capitalize()
    table.cell(0, 3).text                      = 'Hash'.capitalize()
    table.cell(0, 4).text                      = 'Notes'.capitalize()


    # write body cells
    for i in range(1, rows):
        table.cell(i, 0).text                  = df['Repo'][i - 1]
        cell                                   = table.cell(i, 0)
        fill(cell, blue2)
        text_settings(cell, alignment          = PP_ALIGN.CENTER, font_size=Pt(5))
        set_cell_border(cell, blue2, white)


        table.cell(i, 1).text                   = df['Date'][i - 1]
        cell                                    = table.cell(i, 1)
        fill(cell, blue2)
        text_settings(cell, alignment           = PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)

        table.cell(i, 2).text                   = df['Version'][i - 1]
        cell                                    = table.cell(i, 2)
        fill(cell, blue2)
        text_settings(cell, alignment           = PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)

        table.cell(i, 3).text                   = df['Hash'][i - 1]
        cell                                    = table.cell(i, 3)
        fill(cell, blue2)
        text_settings(cell, alignment           = PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)

        table.cell(i, 4).text                   = df['Notes'][i - 1]
        cell                                    = table.cell(i, 4)
        fill(cell, blue2)
        text_settings(cell, alignment           = PP_ALIGN.CENTER)
        set_cell_border(cell, blue2, white)


def set_cell_border(
        cell,
        border_color_LR,
        border_color_TB,
        border_width='12700'):
    """
    Format cell borders.

    :param cell: cell from table
    :param border_color_LR: left and right border colors
    :param border_color_TB: top and bottom border colors
    :param border_width: border width
    """
    # convert RGB to hex
    border_color_LR                             = '%02x%02x%02x' % border_color_LR
    border_color_TB                             = '%02x%02x%02x' % border_color_TB

    colors                                      = [
                                                    border_color_LR,
                                                    border_color_LR,
                                                    border_color_TB,
                                                    border_color_TB]

    tc                                          = cell._tc
    tcPr                                        = tc.get_or_add_tcPr()

    lines                                       = ['a:lnL', 'a:lnR', 'a:lnT', 'a:lnB']

    for line, color in zip(lines, colors):
        ln = SubElement(
            tcPr,
            line,
            w=border_width,
            cap='flat',
            cmpd='sng',
            algn='ctr')
        solidFill                               = SubElement(ln, 'a:solidFill')
        srgbClr                                 = SubElement(solidFill, 'a:srgbClr', val=color)
        prstDash                                = SubElement(ln, 'a:prstDash', val='solid')
        round_                                  = SubElement(ln, 'a:round')
        headEnd                                 = SubElement(ln, 'a:headEnd', type='none', w='med', len='med')
        tailEnd                                 = SubElement(ln, 'a:tailEnd', type='none', w='med', len='med')

def SubElement(parent, tagname, **kwargs):
    element                                     = OxmlElement(tagname)
    element.attrib.update(kwargs)
    parent.append(element)
    return element
